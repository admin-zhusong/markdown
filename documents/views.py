from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.http import JsonResponse, HttpResponse, HttpResponseRedirect
from django.utils import timezone
from django.db.models import Q
from .models import DocumentModel, DocumentVersionModel, ImageModel, FolderModel
import markdown
import os
import uuid
import json
from datetime import timedelta
from django.conf import settings
from io import BytesIO
from xhtml2pdf import pisa

def home_view(request):
    markdown_example = """# 标题一级

## 标题二级

**粗体文本** 和 *斜体文本*

> 这是一段引用

- 无序列表项1
- 无序列表项2

1. 有序列表项1
2. 有序列表项2

```python
def hello():
    print("Hello World")
```

| 姓名 | 年龄 | 城市 |
|------|------|------|
| 张三 | 25 | 北京 |
| 李四 | 30 | 上海 |

[链接文字](https://example.com)

![图片描述](https://via.placeholder.com/150)"""

    markdown_preview = markdown.markdown(markdown_example, extensions=['extra', 'tables', 'fenced_code', 'toc'])

    return render(request, 'documents/home.html', {
        'markdown_example': markdown_example,
        'markdown_preview': markdown_preview
    })

def get_folder_path(folder):
    path = []
    current = folder
    while current:
        path.insert(0, current)
        current = current.parent
    return path

@login_required
def my_documents_view(request):
    search_query = request.GET.get('search', '')
    tag_filter = request.GET.get('tag', '')
    sort_by = request.GET.get('sort', 'update_time')
    sort_order = request.GET.get('order', 'desc')
    folder_id = request.GET.get('folder', None)

    allowed_sort_fields = ['title', 'create_time', 'update_time']
    if sort_by not in allowed_sort_fields:
        sort_by = 'update_time'

    if sort_order not in ['asc', 'desc']:
        sort_order = 'desc'

    current_folder = None
    folder_path = []
    if folder_id and folder_id.isdigit():
        current_folder = get_object_or_404(FolderModel, id=folder_id, user=request.user, is_deleted=False)
        folder_path = get_folder_path(current_folder)

    all_user_documents = DocumentModel.objects.filter(user=request.user, is_deleted=False)
    
    if current_folder:
        all_user_documents = all_user_documents.filter(folder=current_folder)
    else:
        all_user_documents = all_user_documents.filter(folder__isnull=True)

    all_tags = set()
    for doc in all_user_documents:
        for tag in doc.get_tags_list():
            all_tags.add(tag)

    documents = all_user_documents

    if search_query:
        documents = documents.filter(
            Q(title__icontains=search_query) |
            Q(content__icontains=search_query) |
            Q(tags__icontains=search_query)
        )

    if tag_filter:
        documents = documents.filter(tags__icontains=tag_filter)

    sort_field = f'-{sort_by}' if sort_order == 'desc' else sort_by
    documents = documents.order_by(sort_field)

    folders = FolderModel.objects.filter(user=request.user, is_deleted=False, parent=current_folder)

    return render(request, 'documents/my_documents.html', {
        'documents': documents,
        'folders': folders,
        'current_folder': current_folder,
        'folder_path': folder_path,
        'search_query': search_query,
        'tag_filter': tag_filter,
        'all_tags': sorted(all_tags),
        'sort_by': sort_by,
        'sort_order': sort_order,
    })

@login_required
def create_folder_view(request):
    if request.method == 'POST':
        folder_name = request.POST.get('folder_name', '').strip()
        parent_folder_id = request.POST.get('parent_folder', None)
        
        if folder_name:
            parent_folder = None
            if parent_folder_id and parent_folder_id.isdigit():
                parent_folder = get_object_or_404(FolderModel, id=parent_folder_id, user=request.user)
            
            FolderModel.objects.create(
                user=request.user,
                name=folder_name,
                parent=parent_folder
            )
    
    parent_folder_id = request.POST.get('parent_folder', None)
    if parent_folder_id and parent_folder_id.isdigit():
        return HttpResponseRedirect(f'/my-documents/?folder={parent_folder_id}')
    else:
        return HttpResponseRedirect('/my-documents/')

@login_required
def rename_folder_view(request, folder_id):
    folder = get_object_or_404(FolderModel, id=folder_id, user=request.user)
    
    if request.method == 'POST':
        new_name = request.POST.get('folder_name', '').strip()
        if new_name:
            folder.name = new_name
            folder.save()
    
    parent_id = folder.parent.id if folder.parent else None
    if parent_id:
        return HttpResponseRedirect(f'/my-documents/?folder={parent_id}')
    else:
        return HttpResponseRedirect('/my-documents/')

def delete_folder_recursively(folder):
    # 获取所有子文件夹（包括嵌套的）
    child_folders = list(FolderModel.objects.filter(parent=folder, is_deleted=False))
    
    # 递归删除子文件夹
    for child in child_folders:
        delete_folder_recursively(child)
    
    # 删除当前文件夹内的文档
    DocumentModel.objects.filter(folder=folder).update(is_deleted=True, delete_time=timezone.now())
    
    # 删除当前文件夹
    folder.is_deleted = True
    folder.delete_time = timezone.now()
    folder.save()

@login_required
def delete_folder_view(request, folder_id):
    folder = get_object_or_404(FolderModel, id=folder_id, user=request.user)
    parent_folder = folder.parent
    
    if request.method == 'POST':
        try:
            # 递归删除文件夹及其所有内容
            delete_folder_recursively(folder)
        except Exception as e:
            print(f'Error deleting folder: {e}')
    
    parent_id = parent_folder.id if parent_folder else None
    if parent_id:
        return HttpResponseRedirect(f'/my-documents/?folder={parent_id}')
    else:
        return HttpResponseRedirect('/my-documents/')

@login_required
def create_document_view(request):
    folder_id = request.GET.get('folder', None)
    title = f'未命名文档-{timezone.now().strftime("%Y%m%d%H%M%S")}'
    
    folder = None
    if folder_id and folder_id.isdigit():
        folder = get_object_or_404(FolderModel, id=folder_id, user=request.user)
    
    document = DocumentModel.objects.create(
        user=request.user,
        folder=folder,
        title=title,
        content='# 新建文档\n\n开始编写您的Markdown文档吧！',
        original_filename=f'{title}.md'
    )

    DocumentVersionModel.objects.create(
        document=document,
        content=document.content,
        version_num=1,
        remark='初始版本'
    )

    return redirect('edit_document', document_id=document.id)

@login_required
def edit_document_view(request, document_id):
    document = get_object_or_404(DocumentModel, id=document_id, user=request.user)

    if request.method == 'POST':
        document.title = request.POST.get('title', document.title)
        document.content = request.POST.get('content', document.content)
        document.tags = request.POST.get('tags', document.tags)
        document.save()

        DocumentVersionModel.objects.create(
            document=document,
            content=document.content,
            version_num=document.versions.count() + 1,
            remark='手动保存'
        )

        # 保存后重定向回编辑页面，保持在当前文件夹
        return redirect('edit_document', document_id=document.id)

    versions = document.versions.order_by('-create_time')[:10]

    context = {
        'document': document,
        'versions': versions,
    }
    return render(request, 'documents/edit_document.html', context)

@login_required
def document_detail_view(request, document_id):
    document = get_object_or_404(DocumentModel, id=document_id, user=request.user)
    html_content = markdown.markdown(document.content, extensions=['extra', 'tables', 'fenced_code', 'toc'])
    return render(request, 'documents/document_detail.html', {
        'document': document,
        'html_content': html_content
    })

@login_required
def delete_document_view(request, document_id):
    document = get_object_or_404(DocumentModel, id=document_id, user=request.user)

    if request.method == 'POST':
        document.is_deleted = True
        document.delete_time = timezone.now()
        document.save()
        messages.success(request, '文档已删除！')
        return redirect('my_documents')

    return redirect('document_detail', document_id=document_id)

@login_required
def export_document_view(request, document_id):
    document = get_object_or_404(DocumentModel, id=document_id, user=request.user)
    export_format = request.GET.get('format', 'md')

    if export_format == 'md':
        content = document.content
        filename = f'{document.title}.md'
        content_type = 'text/markdown; charset=utf-8'
        response = HttpResponse(content, content_type=content_type)
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response

    elif export_format == 'html':
        html_content = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{document.title}</title>
    <style>
        body {{ font-family: Arial, sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; }}
        h1 {{ border-bottom: 2px solid #333; padding-bottom: 10px; font-size: 24px; }}
        h2 {{ border-bottom: 1px solid #eee; padding-bottom: 5px; font-size: 20px; }}
        h3 {{ font-size: 18px; }}
        code {{ background: #f4f4f4; padding: 2px 6px; border-radius: 3px; font-family: monospace; }}
        pre {{ background: #f4f4f4; padding: 15px; border-radius: 5px; overflow-x: auto; font-family: monospace; }}
        blockquote {{ border-left: 4px solid #ddd; margin-left: 0; padding-left: 15px; color: #666; }}
        table {{ border-collapse: collapse; width: 100%; }}
        th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
        th {{ background: #f4f4f4; }}
        ul, ol {{ padding-left: 20px; }}
        li {{ margin: 5px 0; }}
        p {{ line-height: 1.6; }}
        a {{ color: #007bff; text-decoration: none; }}
        a:hover {{ text-decoration: underline; }}
    </style>
</head>
<body>
{markdown.markdown(document.content, extensions=['extra', 'tables', 'fenced_code', 'toc'])}
</body>
</html>'''
        filename = f'{document.title}.html'
        content_type = 'text/html; charset=utf-8'
        response = HttpResponse(html_content, content_type=content_type)
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response

    elif export_format == 'pdf':
        html_content = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{document.title}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 20px; }}
        h1 {{ border-bottom: 2px solid #333; padding-bottom: 10px; font-size: 22px; }}
        h2 {{ border-bottom: 1px solid #eee; padding-bottom: 5px; font-size: 18px; }}
        h3 {{ font-size: 16px; }}
        code {{ background: #f4f4f4; padding: 2px 6px; border-radius: 3px; font-family: monospace; font-size: 12px; }}
        pre {{ background: #f4f4f4; padding: 12px; border-radius: 4px; font-family: monospace; font-size: 12px; white-space: pre-wrap; word-wrap: break-word; }}
        blockquote {{ border-left: 4px solid #ddd; margin-left: 0; padding-left: 12px; color: #666; }}
        table {{ border-collapse: collapse; width: 100%; }}
        th, td {{ border: 1px solid #ddd; padding: 6px; text-align: left; font-size: 12px; }}
        th {{ background: #f4f4f4; }}
        ul, ol {{ padding-left: 20px; }}
        li {{ margin: 3px 0; font-size: 13px; }}
        p {{ line-height: 1.5; font-size: 13px; }}
        a {{ color: #007bff; }}
    </style>
</head>
<body>
<h1>{document.title}</h1>
{markdown.markdown(document.content, extensions=['extra', 'tables', 'fenced_code', 'toc'])}
</body>
</html>'''

        result = BytesIO()
        pdf = pisa.pisaDocument(BytesIO(html_content.encode('utf-8')), result)

        if pdf.err:
            return HttpResponse('PDF生成失败', status=500)

        response = HttpResponse(result.getvalue(), content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename="{document.title}.pdf"'
        return response

    else:
        content = document.content
        filename = f'{document.title}.md'
        content_type = 'text/markdown; charset=utf-8'
        response = HttpResponse(content, content_type=content_type)
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response

@login_required
def preview_api_view(request):
    if request.method == 'POST':
        content = request.POST.get('content', '')
        html = markdown.markdown(content, extensions=['extra', 'tables', 'fenced_code', 'toc'])
        return HttpResponse(html)
    return HttpResponse('')

@login_required
def autosave_view(request):
    if request.method == 'POST':
        document_id = request.POST.get('document_id')
        content = request.POST.get('content', '')

        try:
            document = DocumentModel.objects.get(id=document_id, user=request.user)
            document.content = content
            document.save()
            return JsonResponse({'status': 'success', 'message': '自动保存成功'})
        except DocumentModel.DoesNotExist:
            return JsonResponse({'status': 'error', 'message': '文档不存在'}, status=404)

    return JsonResponse({'status': 'error', 'message': '无效请求'}, status=400)

def download_view(request, share_token):
    try:
        document = DocumentModel.objects.get(share_token=share_token, is_shared=True)

        if document.share_expire and document.share_expire < timezone.now():
            return HttpResponse('分享链接已过期', status=410)

        if document.share_password:
            if request.session.get('shared_doc_access') != share_token:
                if request.method == 'POST':
                    password = request.POST.get('password', '')
                    if password == document.share_password:
                        request.session['shared_doc_access'] = share_token
                    else:
                        return HttpResponse('密码错误', status=403)
                else:
                    return HttpResponse('需要密码访问', status=403)

        if document.original_filename:
            filename = document.original_filename
            if not filename.lower().endswith('.md'):
                filename = filename + '.md'
        else:
            filename = f"{document.title}.md"

        import urllib.parse
        encoded_filename = urllib.parse.quote(filename)

        content = document.content.encode('utf-8')
        response = HttpResponse(content, content_type='application/octet-stream')
        response['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded_filename}"
        response['Content-Length'] = len(content)
        response['Cache-Control'] = 'no-cache, no-store, must-revalidate'
        response['Pragma'] = 'no-cache'
        response['Expires'] = '0'
        return response

    except DocumentModel.DoesNotExist:
        return HttpResponse('文档不存在或分享已关闭', status=404)

@login_required
def share_api(request, document_id=None):
    if request.method == 'GET':
        if not document_id:
            document_id = request.GET.get('document_id')
        expire_days = int(request.GET.get('expire_days', 0))
    else:
        document_id = request.POST.get('document_id')
        expire_days = int(request.POST.get('expire_days', 0))

    try:
        document = DocumentModel.objects.get(id=document_id, user=request.user)

        if not document.share_token:
            document.share_token = str(uuid.uuid4())[:20]

        if expire_days > 0:
            document.share_expire = timezone.now() + timedelta(days=expire_days)
        else:
            document.share_expire = None

        document.is_shared = True
        document.save()

        share_url = f"{request.scheme}://{request.get_host()}/download/{document.share_token}/"

        return JsonResponse({
            'status': 'success',
            'message': '分享链接已生成',
            'share_token': document.share_token,
            'share_url': share_url
        })

    except DocumentModel.DoesNotExist:
        return JsonResponse({'status': 'error', 'message': '文档不存在'}, status=404)

    return JsonResponse({'status': 'error', 'message': '无效请求'}, status=400)

@login_required
def rename_api(request, document_id):
    if request.method == 'POST':
        title = request.POST.get('title', '').strip()
        
        if not title:
            return JsonResponse({'status': 'error', 'message': '文件名不能为空'}, status=400)

        try:
            document = DocumentModel.objects.get(id=document_id, user=request.user)
            document.title = title
            document.save()
            
            return JsonResponse({'status': 'success', 'message': '重命名成功'})

        except DocumentModel.DoesNotExist:
            return JsonResponse({'status': 'error', 'message': '文档不存在'}, status=404)

    return JsonResponse({'status': 'error', 'message': '无效请求'}, status=400)

def extract_file_content(file, ext):
    """根据文件类型提取内容"""
    if ext in ('.md', '.txt'):
        return file.read().decode('utf-8', errors='ignore')
    
    elif ext in ('.html', '.htm'):
        content = file.read().decode('utf-8', errors='ignore')
        # 简单处理HTML，提取文本内容
        import re
        # 移除HTML标签，保留换行
        text = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
        text = re.sub(r'<[^>]+>', '\n', text)
        # 清理多余的换行
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()
    
    elif ext in ('.docx',):
        try:
            from docx import Document
            doc = Document(file)
            paragraphs = []
            for para in doc.paragraphs:
                paragraphs.append(para.text)
            return '\n\n'.join(paragraphs)
        except Exception as e:
            return f'无法解析Word文档: {str(e)}'
    
    elif ext in ('.xlsx',):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file, read_only=True)
            content = []
            for sheet in wb.worksheets:
                content.append(f'## {sheet.title}')
                for row in sheet.iter_rows(values_only=True):
                    row_content = '\t'.join([str(cell) if cell else '' for cell in row])
                    content.append(row_content)
            return '\n'.join(content)
        except Exception as e:
            return f'无法解析Excel文档: {str(e)}'
    
    elif ext in ('.pptx',):
        try:
            from pptx import Presentation
            prs = Presentation(file)
            content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        content.append(shape.text)
            return '\n\n'.join(content)
        except Exception as e:
            return f'无法解析PPT文档: {str(e)}'
    
    elif ext == '.pdf':
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file)
            content = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    content.append(text)
            return '\n\n'.join(content)
        except Exception as e:
            return f'无法解析PDF文档: {str(e)}'
    
    elif ext == '.doc':
        return '无法直接解析.doc格式，请使用.docx格式或转换后再导入'
    
    elif ext == '.xls':
        return '无法直接解析.xls格式，请使用.xlsx格式或转换后再导入'
    
    elif ext == '.ppt':
        return '无法直接解析.ppt格式，请使用.pptx格式或转换后再导入'
    
    return ''

@login_required
def import_files_view(request):
    if request.method == 'POST':
        csrf_token = request.COOKIES.get('csrftoken')
        import_type = request.POST.get('type', 'single')
        
        # 获取当前文件夹
        folder_id = request.POST.get('folder', None)
        folder = None
        if folder_id and folder_id.isdigit():
            try:
                folder = FolderModel.objects.get(id=folder_id, user=request.user)
            except FolderModel.DoesNotExist:
                folder = None

        if import_type == 'single':
            files = [request.FILES.get('file')]
        else:
            files = request.FILES.getlist('files[]')

        imported_count = 0
        errors = []
        
        # 支持的文件类型
        supported_extensions = {'.md', '.txt', '.html', '.htm', '.docx', '.xlsx', '.pptx', '.pdf', '.doc', '.xls', '.ppt'}

        for file in files:
            if file:
                ext = os.path.splitext(file.name)[1].lower()
                if ext in supported_extensions:
                    try:
                        content = extract_file_content(file, ext)
                        title = os.path.splitext(file.name)[0]
                        DocumentModel.objects.create(
                            user=request.user,
                            title=title,
                            content=content,
                            original_filename=file.name,
                            folder=folder
                        )
                        imported_count += 1
                    except Exception as e:
                        errors.append(f'{file.name}: {str(e)}')
                else:
                    errors.append(f'{file.name}: 不支持的文件类型')

        return JsonResponse({
            'status': 'success',
            'imported': imported_count,
            'errors': errors
        })

    return JsonResponse({'status': 'error', 'message': '无效请求'}, status=400)

@login_required
def version_history_view(request, document_id):
    document = get_object_or_404(DocumentModel, id=document_id, user=request.user)
    versions = document.versions.order_by('-create_time')

    if request.method == 'POST':
        version_id = request.POST.get('version_id')
        version = get_object_or_404(DocumentVersionModel, id=version_id, document=document)
        document.content = version.content
        document.save()
        messages.success(request, f'已恢复到版本 {version.version_num}')
        return redirect('edit_document', document_id=document.id)

    return render(request, 'documents/version_history.html', {
        'document': document,
        'versions': versions
    })

@login_required
def restore_version_view(request, document_id, version_id):
    document = get_object_or_404(DocumentModel, id=document_id, user=request.user)
    version = get_object_or_404(DocumentVersionModel, id=version_id, document=document)

    DocumentVersionModel.objects.create(
        document=document,
        content=document.content,
        version_num=document.versions.count() + 1,
        remark='恢复版本前'
    )

    document.content = version.content
    document.save()

    messages.success(request, f'已恢复到版本 {version.version_num}')
    return redirect('edit_document', document_id=document.id)

@login_required
def recycle_bin_view(request):
    documents = DocumentModel.objects.filter(user=request.user, is_deleted=True).order_by('-delete_time')
    return render(request, 'documents/recycle_bin.html', {'documents': documents})

@login_required
def restore_document_view(request, document_id):
    document = get_object_or_404(DocumentModel, id=document_id, user=request.user)
    document.is_deleted = False
    document.delete_time = None
    document.save()
    messages.success(request, '文档已恢复！')
    return redirect('recycle_bin')

@login_required
def permanent_delete_document_view(request, document_id):
    document = get_object_or_404(DocumentModel, id=document_id, user=request.user)
    document.delete()
    return redirect('recycle_bin')

@login_required
def batch_delete_documents_view(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            document_ids = data.get('document_ids', [])
            
            if document_ids:
                document_ids = [int(doc_id) for doc_id in document_ids]
                DocumentModel.objects.filter(id__in=document_ids, user=request.user, is_deleted=False).update(
                    is_deleted=True,
                    delete_time=timezone.now()
                )
            
            return JsonResponse({'success': True})
        except Exception as e:
            print(f'Error in batch_delete_documents_view: {e}')
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    
    return JsonResponse({'success': False}, status=400)

@login_required
def batch_delete_permanently_view(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            document_ids = data.get('document_ids', [])
            
            if document_ids:
                document_ids = [int(doc_id) for doc_id in document_ids]
                DocumentModel.objects.filter(id__in=document_ids, user=request.user).delete()
            
            return JsonResponse({'success': True})
        except Exception as e:
            print(f'Error in batch_delete_permanently_view: {e}')
            return JsonResponse({'success': False, 'error': str(e)}, status=500)
    
    return JsonResponse({'success': False}, status=400)

def help_view(request):
    return render(request, 'documents/help.html')
